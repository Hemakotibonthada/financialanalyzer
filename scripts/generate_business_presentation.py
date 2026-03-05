"""
FinancialAnalyzer — Premium Business Presentation
Professional enterprise-grade deck with modern design, visual hierarchy,
gradient accents, feature deep-dives, and compelling data storytelling.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ═══════════════════════════════════════════════════════════════
# DESIGN SYSTEM
# ═══════════════════════════════════════════════════════════════

# Primary palette
BG_DARK     = RGBColor(0x0B, 0x0F, 0x1A)
BG_CARD     = RGBColor(0x11, 0x18, 0x27)
BG_CARD2    = RGBColor(0x16, 0x1F, 0x32)
BG_SURFACE  = RGBColor(0x1A, 0x25, 0x3C)

# Accent colors
BLUE        = RGBColor(0x38, 0xBD, 0xF8)
BLUE_DIM    = RGBColor(0x1E, 0x6B, 0x9A)
GREEN       = RGBColor(0x10, 0xB9, 0x81)
GREEN_DIM   = RGBColor(0x06, 0x5F, 0x46)
PURPLE      = RGBColor(0xA7, 0x8B, 0xFA)
PURPLE_DIM  = RGBColor(0x6D, 0x55, 0xC0)
ORANGE      = RGBColor(0xFB, 0xBF, 0x24)
ORANGE_DIM  = RGBColor(0x92, 0x6E, 0x12)
RED         = RGBColor(0xF8, 0x71, 0x71)
RED_DIM     = RGBColor(0x99, 0x33, 0x33)
TEAL        = RGBColor(0x2D, 0xD4, 0xBF)
PINK        = RGBColor(0xF4, 0x72, 0xB6)
YELLOW      = RGBColor(0xFA, 0xCC, 0x15)
CYAN        = RGBColor(0x22, 0xD3, 0xEE)

# Text
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT       = RGBColor(0xE2, 0xE8, 0xF0)
MUTED       = RGBColor(0x94, 0xA3, 0xB8)
DIM         = RGBColor(0x64, 0x74, 0x8B)
FAINT       = RGBColor(0x47, 0x55, 0x69)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ═══════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════

def bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def gradient_bar(slide, left, top, width, height, c1, c2):
    """Add a gradient-filled bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    shape.shadow.inherit = False
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_stops[1].position = 1.0
    return shape

def card(slide, left, top, width, height, color=BG_CARD, radius=Inches(0.15)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def accent_line(slide, left, top, width, color, height=Inches(0.035)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def txt(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, name='Segoe UI'):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return tf

def txt_multi(slide, left, top, width, height, lines, size=14, color=LIGHT, spacing=6):
    """Add multi-paragraph text block."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        p.space_after = Pt(spacing)
    return tf

def stat_box(slide, left, top, w, h, number, label, accent=BLUE, bg_c=BG_CARD):
    c = card(slide, left, top, w, h, bg_c)
    accent_line(slide, left + Inches(0.15), top + Inches(0.12), w - Inches(0.3), accent)
    txt(slide, left, top + Inches(0.3), w, Inches(0.6), str(number), size=36, color=accent, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, left, top + Inches(0.85), w, Inches(0.4), label, size=12, color=MUTED, align=PP_ALIGN.CENTER)

def icon_circle(slide, left, top, icon_char, color, size=Inches(0.5)):
    c = circle(slide, left, top, size, color)
    tf = c.text_frame
    tf.paragraphs[0].text = icon_char
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.bold = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return c

def feature_card(slide, left, top, w, h, icon, title, desc, accent=BLUE, bg_c=BG_CARD):
    c = card(slide, left, top, w, h, bg_c)
    accent_line(slide, left, top, Inches(0.05), accent, h)
    icon_circle(slide, left + Inches(0.2), top + Inches(0.2), icon, accent, Inches(0.45))
    txt(slide, left + Inches(0.8), top + Inches(0.18), w - Inches(1.0), Inches(0.35), title, size=15, color=WHITE, bold=True)
    txt(slide, left + Inches(0.8), top + Inches(0.5), w - Inches(1.0), h - Inches(0.6), desc, size=11, color=MUTED)
    return c

def section_header(slide, number, title, subtitle, accent=BLUE):
    """Consistent section header."""
    bg(slide)
    gradient_bar(slide, Inches(0), Inches(0), W, Inches(0.05), accent, accent)
    # Section number pill
    pill = card(slide, Inches(0.8), Inches(0.5), Inches(0.5), Inches(0.35), accent)
    txt(slide, Inches(0.8), Inches(0.5), Inches(0.5), Inches(0.35), f'{number:02d}', size=14, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, Inches(1.5), Inches(0.45), Inches(4), Inches(0.4), title.upper(), size=13, color=accent, bold=True)
    txt(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.7), subtitle, size=32, color=WHITE, bold=True)
    # Subtle right accent circle
    # Decorative faint circle
    c = circle(slide, Inches(11.5), Inches(0.3), Inches(1.0), BG_SURFACE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE (HERO)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)

# Top gradient strip
gradient_bar(slide, Inches(0), Inches(0), W, Inches(0.08), BLUE, PURPLE)

# Large decorative circles (very faint — use near-background colors)
for cx, cy, sz, clr in [(11.0, 0.5, 2.5, RGBColor(0x10, 0x16, 0x25)), (10.0, 4.5, 3.0, RGBColor(0x12, 0x14, 0x28)), (0.5, 5.5, 2.0, RGBColor(0x0E, 0x16, 0x20))]:
    circle(slide, Inches(cx), Inches(cy), Inches(sz), clr)

# Company badge
pill = card(slide, Inches(0.8), Inches(1.2), Inches(2.8), Inches(0.4), BG_SURFACE)
txt(slide, Inches(0.8), Inches(1.2), Inches(2.8), Inches(0.4), '  CIRCUVENT TECHNOLOGIES', size=11, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# Main title
txt(slide, Inches(0.8), Inches(2.2), Inches(8), Inches(1.2), 'FinancialAnalyzer', size=58, color=WHITE, bold=True)

# Tagline
txt(slide, Inches(0.8), Inches(3.5), Inches(8), Inches(0.6),
    'AI-Powered Enterprise Financial Management Platform', size=22, color=MUTED)

# Subtitle
txt(slide, Inches(0.8), Inches(4.2), Inches(8), Inches(0.5),
    '100% Local ML  •  Indian Market Focus  •  Zero External API Dependencies', size=14, color=DIM)

# Bottom stat cards
stats = [
    ('640+', 'Source Files', BLUE),
    ('150+', 'Features', GREEN),
    ('300+', 'API Endpoints', PURPLE),
    ('33,000+', 'Lines of AI Code', ORANGE),
    ('38', 'ML Modules', TEAL),
    ('0', 'External APIs', RED),
]
for i, (num, label, color) in enumerate(stats):
    x = Inches(0.8 + i * 2.05)
    stat_box(slide, x, Inches(5.4), Inches(1.85), Inches(1.3), num, label, color)

# Version footer
txt(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
    'Version 2.0.0  |  March 2026  |  Business Presentation', size=11, color=FAINT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM WE SOLVE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 0, 'The Challenge', 'Why FinancialAnalyzer?')

problems = [
    ('⚡', 'Manual Tracking', 'Spreadsheets and notebooks can\'t scale. Users lose\nhours tracking expenses, EMIs, and investments manually.', RED),
    ('🔒', 'Data Privacy', 'Cloud-based AI tools send all your financial data\nto external servers. FinancialAnalyzer keeps it 100% local.', ORANGE),
    ('🧩', 'Fragmented Tools', 'Separate apps for budgets, investments, taxes, and\ndebt. One platform unifies everything.', PURPLE),
    ('🧠', 'No Intelligence', 'Most finance apps just record numbers. They don\'t\npredict, optimize, or detect anomalies.', BLUE),
    ('🇮🇳', 'Not India-Ready', 'Most fintech tools target US markets. We support\nINR, CIBIL, Indian tax rules, UPI, NPS, PPF, ELSS.', GREEN),
    ('💰', 'Expensive AI', 'GPT/Claude API costs add up. Our 38 ML modules\nrun entirely on your server — zero API fees.', TEAL),
]

for i, (icon, title, desc, color) in enumerate(problems):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(2.0 + row * 2.5)
    feature_card(slide, x, y, Inches(3.8), Inches(2.1), icon[0], title, desc, color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: PLATFORM OVERVIEW
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 1, 'Platform Overview', 'Full-Stack Architecture')

# Architecture layers with gradient cards
layers = [
    ('React 18 + Vite 5', 'SPA with 146 pages, MUI 7, Tailwind CSS\n3-mode theme (Light / Dark / AMOLED)', BLUE, '🖥️'),
    ('Express.js API', '70+ route groups, 300+ endpoints\nWebSocket real-time, Enterprise middleware', GREEN, '⚙️'),
    ('AI/ML Engine', '38 modules, 33K lines, Neural Nets\nRL, NLP, Anomaly Detection — all from scratch', PURPLE, '🧠'),
    ('MongoDB + Redis', '47 models, per-user ML persistence\nTTL indexes, file-based model storage', ORANGE, '🗄️'),
]

for i, (title, desc, color, icon) in enumerate(layers):
    y = Inches(1.9 + i * 1.25)
    # Layer card
    c = card(slide, Inches(0.8), y, Inches(7.5), Inches(1.1), BG_CARD)
    # Left accent border
    accent_line(slide, Inches(0.8), y, Inches(0.06), color, Inches(1.1))
    # Icon
    icon_circle(slide, Inches(1.1), y + Inches(0.3), icon[0], color, Inches(0.5))
    # Text
    txt(slide, Inches(1.8), y + Inches(0.1), Inches(3.0), Inches(0.35), title, size=17, color=color, bold=True)
    txt(slide, Inches(1.8), y + Inches(0.45), Inches(5.5), Inches(0.6), desc, size=12, color=MUTED)
    # Connection line
    if i < len(layers) - 1:
        accent_line(slide, Inches(1.35), y + Inches(1.1), Inches(0.02), DIM, Inches(0.15))

# Right side — key metrics
card(slide, Inches(9.0), Inches(1.9), Inches(3.5), Inches(5.0), BG_CARD)
accent_line(slide, Inches(9.0), Inches(1.9), Inches(3.5), BLUE)
txt(slide, Inches(9.2), Inches(2.1), Inches(3.0), Inches(0.35), 'Key Metrics', size=16, color=BLUE, bold=True)

metrics = [
    ('89', 'Backend Route Files'),
    ('128', 'Service Modules'),
    ('47', 'Database Models'),
    ('146', 'Frontend Pages'),
    ('67', 'UI Components'),
    ('8', 'React Contexts'),
    ('12', 'Middleware Layers'),
    ('25+', 'ML Algorithms'),
]

for i, (num, label) in enumerate(metrics):
    y = Inches(2.6 + i * 0.52)
    txt(slide, Inches(9.4), y, Inches(0.7), Inches(0.3), num, size=16, color=BLUE, bold=True, align=PP_ALIGN.RIGHT)
    txt(slide, Inches(10.2), y, Inches(2.0), Inches(0.3), label, size=12, color=MUTED)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: CORE — DASHBOARD & HEALTH
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 2, 'Core Platform', 'Dashboards & Financial Health')

# Dashboard versions
versions = [
    ('Dashboard V1', 'income/expense summary, recent\ntransactions, budget progress,\nmonthly trends chart', BLUE, 'Functional'),
    ('Dashboard V2', 'Animated KPIs, spending trend\ncharts, category intelligence,\nAI predictions panel', GREEN, 'Enhanced'),
    ('Dashboard V3', 'Enterprise-grade, real-time\nmetrics, department analytics,\nAI-powered forecasts', PURPLE, 'Enterprise'),
]

for i, (title, desc, color, badge) in enumerate(versions):
    x = Inches(0.8 + i * 4.1)
    c = card(slide, x, Inches(2.0), Inches(3.8), Inches(2.6), BG_CARD)
    accent_line(slide, x, Inches(2.0), Inches(3.8), color)
    # Version badge
    pill_card = card(slide, x + Inches(0.15), Inches(2.15), Inches(1.2), Inches(0.3), color)
    txt(slide, x + Inches(0.15), Inches(2.15), Inches(1.2), Inches(0.3), badge, size=10, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.2), Inches(2.55), Inches(3.4), Inches(0.35), title, size=17, color=WHITE, bold=True)
    txt(slide, x + Inches(0.2), Inches(2.95), Inches(3.4), Inches(1.5), desc, size=12, color=MUTED)

# Health Score section
card(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(2.0), BG_CARD)
accent_line(slide, Inches(0.8), Inches(5.0), Inches(11.7), GREEN)

# Score circle mock
score_circle = circle(slide, Inches(1.2), Inches(5.3), Inches(1.4), GREEN_DIM)
txt(slide, Inches(1.2), Inches(5.5), Inches(1.4), Inches(0.6), '78', size=36, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
txt(slide, Inches(1.2), Inches(6.1), Inches(1.4), Inches(0.3), 'SCORE', size=10, color=MUTED, align=PP_ALIGN.CENTER)

txt(slide, Inches(3.0), Inches(5.2), Inches(3.5), Inches(0.35), 'Financial Health Score', size=18, color=WHITE, bold=True)
txt(slide, Inches(3.0), Inches(5.6), Inches(4.0), Inches(1.2),
    '0-100 composite score across savings rate,\ndebt-to-income ratio, emergency fund coverage,\ninsurance adequacy, and investment diversification.\n3 versions: V1 → V2 → Enterprise V3', size=12, color=MUTED)

# Health dimensions
dims = ['Savings', 'Debt', 'Emergency', 'Insurance', 'Investment', 'Goals']
colors_dim = [GREEN, RED, ORANGE, BLUE, PURPLE, TEAL]
for i, (dim, clr) in enumerate(zip(dims, colors_dim)):
    x = Inches(7.5 + (i % 3) * 1.5)
    y = Inches(5.3 + (i // 3) * 0.8)
    icon_circle(slide, x, y + Inches(0.05), '●', clr, Inches(0.3))
    txt(slide, x + Inches(0.4), y, Inches(1.0), Inches(0.35), dim, size=11, color=LIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: TRANSACTION MANAGEMENT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 3, 'Transactions', 'Smart Transaction Management')

# Input sources
sources = [
    ('📝', 'Manual Entry', 'Quick add via navbar or\ndashboard widget', BLUE),
    ('📄', 'Document Upload', 'PDF / CSV / image bank\nstatement parsing', GREEN),
    ('📧', 'Gmail Sync', 'Auto-extract from bank\nalert emails (UPI, NEFT)', PURPLE),
    ('📸', 'Receipt Scanner', 'Tesseract OCR with\nauto-categorization', ORANGE),
    ('📊', 'CSV Import', 'Bulk import with\ncolumn mapping', TEAL),
    ('💳', 'UPI Detection', 'Extract VPA, UTR,\napp details auto', PINK),
]

for i, (icon, title, desc, color) in enumerate(sources):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 2.0)
    y = Inches(2.0 + row * 2.0)
    c = card(slide, x, y, Inches(1.8), Inches(1.7), BG_CARD)
    accent_line(slide, x, y, Inches(1.8), color)
    txt(slide, x, y + Inches(0.2), Inches(1.8), Inches(0.3), icon, size=20, align=PP_ALIGN.CENTER)
    txt(slide, x, y + Inches(0.55), Inches(1.8), Inches(0.3), title, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x, y + Inches(0.9), Inches(1.8), Inches(0.7), desc, size=10, color=MUTED, align=PP_ALIGN.CENTER)

# AI Pipeline (right side)
card(slide, Inches(7.2), Inches(2.0), Inches(5.3), Inches(4.8), BG_CARD)
gradient_bar(slide, Inches(7.2), Inches(2.0), Inches(5.3), Inches(0.04), PURPLE, BLUE)
txt(slide, Inches(7.5), Inches(2.2), Inches(4.8), Inches(0.4), 'AI Enrichment Pipeline', size=18, color=PURPLE, bold=True)

pipeline_steps = [
    ('1', 'Category Classification', 'ML model trained on user data', BLUE),
    ('2', 'Merchant Normalization', 'Standardize merchant names', GREEN),
    ('3', 'Anomaly Scoring', 'Isolation Forest detection', RED),
    ('4', 'UPI Metadata Extract', 'VPA, UTR, app details', PURPLE),
    ('5', 'Recurring Detection', 'Identify subscriptions', ORANGE),
    ('6', 'Tax Deductibility', 'Auto-flag tax-saving txns', TEAL),
    ('7', 'Budget Impact', 'Real-time budget update', BLUE),
    ('8', 'Alert Triggers', 'Budget warnings, anomalies', RED),
]

for i, (num, title, desc, color) in enumerate(pipeline_steps):
    y = Inches(2.8 + i * 0.48)
    # Step number
    pill_card = card(slide, Inches(7.5), y, Inches(0.35), Inches(0.3), color)
    txt(slide, Inches(7.5), y, Inches(0.35), Inches(0.3), num, size=10, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, Inches(8.0), y, Inches(2.2), Inches(0.3), title, size=12, color=WHITE, bold=True)
    txt(slide, Inches(10.3), y, Inches(2.0), Inches(0.3), desc, size=10, color=DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: BUDGETS & PLANNING
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 4, 'Budgeting', 'Smart Budget Management & Planning')

# Budget features grid
budget_features = [
    ('📋', 'Category Budgets', 'Monthly/weekly/yearly budgets\nper expense category with\nreal-time progress tracking', BLUE),
    ('✨', '50/30/20 Auto-Fill', 'Auto-calculate budgets from\nsalary: 50% needs, 30% wants,\n20% savings/investments', GREEN),
    ('🤖', 'AI Optimizer', 'Reinforcement Learning finds\noptimal allocation across\ncategories for your goals', PURPLE),
    ('⚠️', 'Smart Alerts', 'Configurable thresholds (80%\ndefault), multi-channel alerts\n(push, email, SMS)', ORANGE),
    ('🔄', 'Budget Rollover', 'Carry forward unused budget\nto next period automatically\nwith rollover tracking', TEAL),
    ('📅', 'Financial Calendar', 'Calendar view of all upcoming\nbills, EMIs, subscriptions,\nand financial events', PINK),
    ('🧙', 'Smart Wizard', 'Step-by-step guided budget\ncreation with AI suggestions\nbased on spending patterns', CYAN),
    ('📊', 'Budget Intelligence', 'Enterprise V3 ML-driven\nbudget recommendations,\ntrend analysis & forecasting', YELLOW),
]

for i, (icon, title, desc, color) in enumerate(budget_features):
    col = i % 4
    row = i // 4
    x = Inches(0.8 + col * 3.1)
    y = Inches(2.0 + row * 2.5)
    feature_card(slide, x, y, Inches(2.85), Inches(2.1), icon[0], title, desc, color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: INVESTMENTS & MARKETS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 5, 'Investments', 'Portfolio Management & Market Intelligence')

# Investment types (left)
card(slide, Inches(0.8), Inches(2.0), Inches(4.0), Inches(4.8), BG_CARD)
gradient_bar(slide, Inches(0.8), Inches(2.0), Inches(4.0), Inches(0.04), BLUE, GREEN)
txt(slide, Inches(1.0), Inches(2.2), Inches(3.5), Inches(0.4), 'Supported Investment Types', size=16, color=BLUE, bold=True)

inv_types = [
    ('Stocks & Equities', BLUE), ('Mutual Funds', GREEN), ('Fixed Deposits', ORANGE),
    ('Gold (Physical + Digital)', YELLOW), ('Cryptocurrency', PURPLE), ('PPF / EPF / NPS', TEAL),
    ('ELSS (Tax Saving)', GREEN), ('REITs & Bonds', BLUE), ('SIPs', CYAN),
]

for i, (inv, color) in enumerate(inv_types):
    y = Inches(2.8 + i * 0.42)
    circle(slide, Inches(1.2), y + Inches(0.08), Inches(0.2), color)
    txt(slide, Inches(1.55), y, Inches(3.0), Inches(0.3), inv, size=12, color=LIGHT)

# Analytics (right)
card(slide, Inches(5.2), Inches(2.0), Inches(7.3), Inches(4.8), BG_CARD)
gradient_bar(slide, Inches(5.2), Inches(2.0), Inches(7.3), Inches(0.04), PURPLE, ORANGE)
txt(slide, Inches(5.5), Inches(2.2), Inches(6.8), Inches(0.4), '22 Analytics Endpoints', size=16, color=PURPLE, bold=True)

analytics = [
    ('Portfolio Dashboard & Summary', 'XIRR, CAGR, Absolute Returns'),
    ('Asset Allocation Analysis', 'Visual allocation across types'),
    ('Risk Analysis & VaR', 'Risk scoring, Value at Risk'),
    ('Diversification Score', '0-100 portfolio health score'),
    ('Tax Efficiency Analysis', 'STCG / LTCG impact analysis'),
    ('AI Recommendations', 'Buy / Sell / Hold signals'),
    ('Markowitz Optimization', 'Mean-variance efficient frontier'),
    ('Monte Carlo Simulation', 'Scenario analysis with intervals'),
    ('Black-Litterman Model', 'Incorporating investor views'),
]

for i, (title, desc) in enumerate(analytics):
    y = Inches(2.8 + i * 0.44)
    txt(slide, Inches(5.5), y, Inches(3.5), Inches(0.3), f'▸  {title}', size=12, color=WHITE, bold=True)
    txt(slide, Inches(9.2), y, Inches(3.0), Inches(0.3), desc, size=10, color=DIM)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: DEBT, EMI & LOANS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 6, 'Debt & EMI', 'Comprehensive Debt Freedom System')

# Big stats
stat_box(slide, Inches(0.8), Inches(1.85), Inches(2.8), Inches(1.3), '33', 'EMI API Endpoints', RED, BG_CARD)
stat_box(slide, Inches(3.85), Inches(1.85), Inches(2.8), Inches(1.3), '4,384', 'Lines (EMI Routes)', ORANGE, BG_CARD)
stat_box(slide, Inches(6.9), Inches(1.85), Inches(2.8), Inches(1.3), '3', 'UI Versions', PURPLE, BG_CARD)
stat_box(slide, Inches(9.95), Inches(1.85), Inches(2.8), Inches(1.3), '8', 'Debt Types', BLUE, BG_CARD)

# Feature columns
col1 = [
    ('One-Click Prepay', 'Create prepayment intent with\nauto surplus fund sweep'),
    ('Balance Transfer', 'Analyze transfer opportunities\nfor lower interest rates'),
    ('Foreclosure Calc', 'Calculate early closure costs\nand total savings'),
    ('PDF/Excel Export', 'Export amortization schedules\nin multiple formats'),
]

col2 = [
    ('Late Fee Shield', 'Automated late fee protection\nwith pre-due reminders'),
    ('Snowball vs Avalanche', 'Compare payoff strategies\n+ AI hybrid recommendation'),
    ('Bank Deduction', 'Track auto-deduction from\nlinked bank accounts'),
    ('Emergency Fund', 'Track fund with contribution\nhistory and goal progress'),
]

for i, (title, desc) in enumerate(col1):
    y = Inches(3.5 + i * 0.95)
    feature_card(slide, Inches(0.8), y, Inches(5.7), Inches(0.8), '▸', title, desc, RED, BG_CARD)

for i, (title, desc) in enumerate(col2):
    y = Inches(3.5 + i * 0.95)
    feature_card(slide, Inches(6.8), y, Inches(5.7), Inches(0.8), '▸', title, desc, ORANGE, BG_CARD)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: AI ENGINE OVERVIEW
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 7, 'AI Engine', 'Machine Learning — 100% Local, Zero APIs')

# Hero statement
card(slide, Inches(0.8), Inches(1.85), Inches(11.7), Inches(1.0), BG_SURFACE)
txt(slide, Inches(1.0), Inches(1.95), Inches(11.3), Inches(0.35),
    'Every algorithm implemented from scratch in JavaScript — no TensorFlow, no PyTorch, no OpenAI.',
    size=16, color=LIGHT, bold=True, align=PP_ALIGN.CENTER)
txt(slide, Inches(1.0), Inches(2.35), Inches(11.3), Inches(0.35),
    'Full data privacy. Zero API costs. Unlimited predictions. Works offline.',
    size=13, color=MUTED, align=PP_ALIGN.CENTER)

# Category cards
ai_categories = [
    ('🔧', 'Core ML', '4,650 lines', 'Neural Networks, Decision Trees,\nRandom Forest, K-Means, PCA,\nARIMA, Holt-Winters', BLUE),
    ('🎯', 'Training', '4,880 lines', 'Auto Pipeline, Self-Learning,\nAutoML, Model Registry,\nDrift Detection, A/B Testing', GREEN),
    ('📈', 'Prediction', '2,360 lines', 'Financial Forecasting, CIBIL\nScore Prediction, Cash Flow\nIntelligence & Projection', CYAN),
    ('🚨', 'Anomaly', '2,510 lines', 'Isolation Forest, LOF,\nAutoencoder, SPC Control\nCharts, Fraud Detection', RED),
    ('💬', 'NLP', '4,270 lines', 'TF-IDF, NER, Sentiment,\nConversational AI, Semantic\nSearch, NL Report Gen', PURPLE),
    ('⚡', 'Optimization', '3,970 lines', 'RL (DQN, Actor-Critic),\nPortfolio (Markowitz),\nGoals, Tax Harvesting', ORANGE),
    ('🧠', 'Behavioral', '3,200 lines', 'Bias Detection, Spending\nPatterns, Peer Comparison,\nFinancial Wellness Score', TEAL),
    ('🔍', 'Explainable', '1,805 lines', 'SHAP, LIME, Counterfactual\nAnalysis, Knowledge Graph,\nDecision Audit Trail', PINK),
]

for i, (icon, title, lines, desc, color) in enumerate(ai_categories):
    col = i % 4
    row = i // 4
    x = Inches(0.8 + col * 3.1)
    y = Inches(3.2 + row * 2.1)
    c = card(slide, x, y, Inches(2.85), Inches(1.8), BG_CARD)
    accent_line(slide, x, y, Inches(2.85), color)
    txt(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.3), icon, size=18, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.55), y + Inches(0.12), Inches(1.5), Inches(0.3), title, size=14, color=WHITE, bold=True)
    # Lines badge
    pill_card = card(slide, x + Inches(1.7), y + Inches(0.15), Inches(1.0), Inches(0.25), color)
    txt(slide, x + Inches(1.7), y + Inches(0.15), Inches(1.0), Inches(0.25), lines, size=9, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.15), y + Inches(0.5), Inches(2.55), Inches(1.2), desc, size=11, color=MUTED)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: AI LAB — 13 DEDICATED PAGES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 8, 'AI Lab', '13 Dedicated AI Feature Pages')

ai_lab_pages = [
    ('💬', 'AI Chatbot', 'Multi-turn conversational\nfinancial AI assistant', BLUE),
    ('🎮', 'RL Optimizer', 'Q-Learning / DQN budget\n& investment optimizer', GREEN),
    ('📡', 'Model Observatory', 'ML performance monitoring\n& drift detection', PURPLE),
    ('🚨', 'Anomaly Lab', 'Multi-algorithm ensemble\ndetection (IF, LOF, SPC)', RED),
    ('📋', 'Financial Planner', 'AI-powered comprehensive\nlifetime financial plan', ORANGE),
    ('🔍', 'Spending Intel', 'Merchant-level analysis,\nimpulse spend detection', TEAL),
    ('📊', 'Portfolio AI', 'Markowitz, Black-Litterman,\nEfficient Frontier gen', BLUE),
    ('💳', 'Credit Predictor', 'CIBIL score prediction +\nwhat-if simulation', PINK),
    ('💰', 'Cash Flow AI', 'Income pattern detection,\nliquidity assessment', GREEN),
    ('🔔', 'Sub Manager', 'Auto-detect 55+ Indian\nsubscriptions, optimize', YELLOW),
    ('🎯', 'Goal & Tax AI', 'Goal achievement + tax\nharvesting (STCG/LTCG)', CYAN),
    ('❤️', 'Wellness AI', '8-dimension financial\nwellness scoring', RED),
    ('🏠', 'AI Command V3', 'Central control panel\nwith health monitoring', PURPLE),
]

for i, (icon, title, desc, color) in enumerate(ai_lab_pages):
    if i < 7:
        col = i
        row = 0
    else:
        col = i - 7
        row = 1
    
    x = Inches(0.3 + col * 1.82)
    y = Inches(2.0 + row * 2.5)
    
    c = card(slide, x, y, Inches(1.65), Inches(2.1), BG_CARD)
    accent_line(slide, x, y, Inches(1.65), color)
    # Icon centered
    txt(slide, x, y + Inches(0.2), Inches(1.65), Inches(0.4), icon, size=24, align=PP_ALIGN.CENTER)
    txt(slide, x, y + Inches(0.65), Inches(1.65), Inches(0.3), title, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.1), y + Inches(1.0), Inches(1.45), Inches(0.9), desc, size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11: NLP & CONVERSATIONAL AI
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 9, 'NLP & Chat', 'Natural Language Processing & AI Conversations')

# NLP capabilities (left)
card(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(5.0), BG_CARD)
gradient_bar(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(0.04), PURPLE, BLUE)
txt(slide, Inches(1.0), Inches(2.2), Inches(5.4), Inches(0.35), 'NLP Capabilities (5 modules, 4,270 lines)', size=15, color=PURPLE, bold=True)

nlp_feats = [
    ('Tokenizer', 'Porter-like stemming, financial term preservation'),
    ('TF-IDF', 'Document vectorization & similarity matching'),
    ('Sentiment', 'Transaction-level sentiment analysis'),
    ('Financial NER', 'Extract ₹ amounts, dates, accounts, merchants'),
    ('Query Engine', '"Show me food spending > ₹500 last month"'),
    ('Summarizer', 'Auto-generate financial narrative reports'),
    ('Semantic Search', 'Fuzzy matching, intent-aware, faceted search'),
    ('Doc Intelligence', 'Parse bank statements, salary slips, tax docs'),
]

for i, (title, desc) in enumerate(nlp_feats):
    y = Inches(2.7 + i * 0.5)
    icon_circle(slide, Inches(1.1), y + Inches(0.05), '▸', PURPLE, Inches(0.25))
    txt(slide, Inches(1.5), y, Inches(2.0), Inches(0.3), title, size=12, color=WHITE, bold=True)
    txt(slide, Inches(3.5), y, Inches(2.8), Inches(0.3), desc, size=10, color=MUTED)

# Conversational AI (right)
card(slide, Inches(6.9), Inches(2.0), Inches(5.6), Inches(5.0), BG_CARD)
gradient_bar(slide, Inches(6.9), Inches(2.0), Inches(5.6), Inches(0.04), GREEN, TEAL)
txt(slide, Inches(7.1), Inches(2.2), Inches(5.2), Inches(0.35), 'Conversational AI (1,291 lines)', size=15, color=GREEN, bold=True)

chat_feats = [
    'Multi-turn dialogue with memory persistence',
    'Short-term + long-term conversation memory',
    'Intent classification (40+ financial intents)',
    'Entity extraction (amounts, dates, categories)',
    'Context tracking across conversation turns',
    'Preference learning from interactions',
    '3 chat versions (V1, V2, Enterprise V3)',
    '',
    'Example queries:',
    '"How much did I spend on Swiggy?"',
    '"Compare this month vs last month"',
    '"What\'s my savings rate this quarter?"',
]

for i, line in enumerate(chat_feats):
    y = Inches(2.7 + i * 0.38)
    color = GREEN if line.startswith('"') else (WHITE if line == 'Example queries:' else MUTED)
    bold = line == 'Example queries:' or line.startswith('"')
    txt(slide, Inches(7.3), y, Inches(5.0), Inches(0.3), line if not line.startswith('"') else f'  {line}', size=11, color=color, bold=bold)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12: ANOMALY DETECTION & FRAUD
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 10, 'Security', 'Anomaly Detection & Fraud Prevention')

# Anomaly detection algorithms
card(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(5.0), BG_CARD)
gradient_bar(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(0.04), RED, ORANGE)
txt(slide, Inches(1.0), Inches(2.2), Inches(5.4), Inches(0.35), 'Anomaly Detection — 5 Ensemble Algorithms', size=15, color=RED, bold=True)

algos = [
    ('1', 'Isolation Forest', 'Random partitioning — anomalies\nisolate faster than normal data', RED),
    ('2', 'Local Outlier Factor', 'Local density comparison —\nfinds contextual outliers', ORANGE),
    ('3', 'Statistical Process Control', 'Shewhart control charts with\nupper/lower control limits', YELLOW),
    ('4', 'Autoencoder', 'Neural network reconstruction\nerror — high error = anomaly', PURPLE),
    ('5', 'DBSCAN Clustering', 'Density-based outlier detection\nfor group-level anomalies', BLUE),
]

for i, (num, title, desc, color) in enumerate(algos):
    y = Inches(2.8 + i * 0.85)
    pill_card = card(slide, Inches(1.1), y, Inches(0.3), Inches(0.3), color)
    txt(slide, Inches(1.1), y, Inches(0.3), Inches(0.3), num, size=12, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, Inches(1.6), y, Inches(2.5), Inches(0.3), title, size=13, color=WHITE, bold=True)
    txt(slide, Inches(1.6), y + Inches(0.3), Inches(4.5), Inches(0.4), desc, size=10, color=MUTED)

# Fraud detection layers
card(slide, Inches(6.9), Inches(2.0), Inches(5.6), Inches(5.0), BG_CARD)
gradient_bar(slide, Inches(6.9), Inches(2.0), Inches(5.6), Inches(0.04), ORANGE, RED)
txt(slide, Inches(7.1), Inches(2.2), Inches(5.2), Inches(0.35), 'Fraud Prevention — 5 Defense Layers', size=15, color=ORANGE, bold=True)

fraud_layers = [
    ('🛡️', 'Rule Engine', 'Configurable rules: high-value\nthreshold > 5× average, velocity\nchecks, daily spending limits', RED),
    ('🧠', 'ML Scoring', 'Anomaly detection ensemble\nintegration for real-time\ntransaction scoring', ORANGE),
    ('👤', 'Behavioral', 'Spending pattern deviation\nanalysis detects compromised\naccount activity', YELLOW),
    ('⚡', 'Velocity', 'Transaction frequency + amount\nvelocity checks per timeframe', PURPLE),
    ('📍', 'Geolocation', 'Location-based anomaly\ndetection for unusual\ntransaction origins', BLUE),
]

for i, (icon, title, desc, color) in enumerate(fraud_layers):
    y = Inches(2.75 + i * 0.82)
    feature_card(slide, Inches(7.1), y, Inches(5.2), Inches(0.7), icon[0], title, desc, color, BG_SURFACE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13: ENTERPRISE FEATURES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 11, 'Enterprise', '3-Generation Progressive UI Architecture')

# V1 / V2 / V3 comparison
gens = [
    ('V1', 'Functional', '80+ pages', 'Core financial operations\nTailwind-based UI\nFull feature coverage', BLUE, '✅ Stable'),
    ('V2', 'Enhanced', '11 pages', 'MUI 7 components\nAnimated KPIs & charts\nEnhanced UX patterns', GREEN, '🚀 Enhanced'),
    ('V3', 'Enterprise', '13 pages', 'Real-time analytics\nAI-powered dashboards\nDepartment-level insights', PURPLE, '⚡ Enterprise'),
]

for i, (ver, label, pages, desc, color, badge) in enumerate(gens):
    x = Inches(0.8 + i * 4.1)
    c = card(slide, x, Inches(2.0), Inches(3.8), Inches(2.8), BG_CARD)
    accent_line(slide, x, Inches(2.0), Inches(3.8), color)
    # Version badge
    pill_card = card(slide, x + Inches(0.15), Inches(2.15), Inches(0.5), Inches(0.35), color)
    txt(slide, x + Inches(0.15), Inches(2.15), Inches(0.5), Inches(0.35), ver, size=14, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.8), Inches(2.17), Inches(1.5), Inches(0.3), label, size=14, color=WHITE, bold=True)
    pill_card2 = card(slide, x + Inches(2.5), Inches(2.2), Inches(1.1), Inches(0.22), BG_SURFACE)
    txt(slide, x + Inches(2.5), Inches(2.2), Inches(1.1), Inches(0.22), pages, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.2), Inches(2.65), Inches(3.4), Inches(1.5), desc, size=12, color=MUTED)
    txt(slide, x + Inches(0.2), Inches(4.1), Inches(3.4), Inches(0.3), badge, size=12, color=color, bold=True)

# Enterprise-only features
enterprise_feats = [
    ('Admin Panel', 'User management, system health, role control', RED),
    ('Automation Engine', 'If-then automation rules with cron scheduling', PURPLE),
    ('Smart Notifications', 'AI priority P0-P3, fatigue prevention, timing', BLUE),
    ('Activity Logging', 'Full audit trail — every API request logged', GREEN),
    ('Enterprise Reports', 'Scheduled PDF/Excel delivery with branding', ORANGE),
    ('Lender Dashboard', 'Role-based dashboard for money lenders', TEAL),
]

for i, (title, desc, color) in enumerate(enterprise_feats):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(5.2 + row * 0.95)
    feature_card(slide, x, y, Inches(3.8), Inches(0.8), '▸', title, desc, color, BG_CARD)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14: REPORTS, EXPORT & ANALYTICS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 12, 'Reports', 'Analytics, Reports & Data Export')

# Report features
report_feats = [
    ('📊', 'Reports Hub', 'Financial report generation\nwith 3 UI versions', BLUE),
    ('📈', 'Advanced Analytics', 'Multi-dimensional analysis\nwith category drill-down', GREEN),
    ('🔬', 'Data Viz Lab', 'Interactive data visualization\nplayground', PURPLE),
    ('📋', 'Export Center', 'PDF / Excel / CSV output\nwith custom templates', ORANGE),
    ('⚖️', 'Comparison Tool', 'Month-over-month, year-over-\nyear analysis', TEAL),
    ('🏆', 'Scorecard', 'Comprehensive financial\nscoring across all metrics', PINK),
    ('📉', 'Risk Dashboard', 'Financial risk assessment\nwith drill-down', RED),
    ('🔍', 'Spending Insights', 'Deep spending analysis\n(2 versions)', CYAN),
]

for i, (icon, title, desc, color) in enumerate(report_feats):
    col = i % 4
    row = i // 4
    x = Inches(0.8 + col * 3.1)
    y = Inches(2.0 + row * 2.2)
    c = card(slide, x, y, Inches(2.85), Inches(1.8), BG_CARD)
    accent_line(slide, x, y, Inches(2.85), color)
    txt(slide, x, y + Inches(0.2), Inches(2.85), Inches(0.4), icon, size=22, align=PP_ALIGN.CENTER)
    txt(slide, x, y + Inches(0.55), Inches(2.85), Inches(0.3), title, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.15), y + Inches(0.9), Inches(2.55), Inches(0.7), desc, size=11, color=MUTED, align=PP_ALIGN.CENTER)

# Export format bar
card(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6), BG_SURFACE)
formats = [('PDF', '📄', BLUE), ('Excel', '📗', GREEN), ('CSV', '📊', ORANGE), ('JSON', '{}', PURPLE)]
for i, (fmt, icon, color) in enumerate(formats):
    x = Inches(2.0 + i * 2.8)
    icon_circle(slide, x, Inches(6.65), icon[0], color, Inches(0.4))
    txt(slide, x + Inches(0.5), Inches(6.7), Inches(1.5), Inches(0.4), fmt, size=14, color=WHITE, bold=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 15: INTEGRATIONS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 13, 'Integrations', 'Platform Integrations & Connectivity')

integrations = [
    ('📧', 'Gmail API', 'Auto-sync bank email alerts\n(UPI, NEFT, CC, RTGS)\n1-year lookback on first sync', BLUE),
    ('☁️', 'Google Drive', 'Cloud backup and restore\nof all financial data with\nscheduled auto-sync', GREEN),
    ('🏦', 'Plaid Banking', 'Direct bank connection\ntransaction import, balance\nsynchronization', PURPLE),
    ('🔐', 'Firebase Auth', 'Google OAuth login +\ncloud hosting deployment\nDual auth backend', ORANGE),
    ('📱', 'Twilio SMS', 'SMS notifications for\ncritical alerts: fraud,\nEMI due, budget exceeded', RED),
    ('⚡', 'WebSocket', 'Real-time notifications,\nlive dashboard updates,\ncross-device sync', TEAL),
    ('📸', 'Tesseract OCR', 'Receipt & document text\nextraction for auto\ntransaction creation', PINK),
    ('🖥️', 'Electron App', 'Native Windows/Mac\ndesktop application\nwith auto-update', CYAN),
    ('📱', 'PWA Support', 'Progressive Web App +\noffline capability\nwith service worker', YELLOW),
]

for i, (icon, title, desc, color) in enumerate(integrations):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(2.0 + row * 1.7)
    feature_card(slide, x, y, Inches(3.8), Inches(1.4), icon[0], title, desc, color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 16: SECURITY & AUTH
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 14, 'Security', 'Enterprise-Grade Security Architecture')

security_layers = [
    ('🔐', 'JWT Authentication', 'Access + Refresh token pair with\nautomatic rotation and revocation', RED),
    ('🔑', 'Two-Factor Auth', 'TOTP-based 2FA via Google\nAuthenticator or Authy app', ORANGE),
    ('🔒', 'Password Security', 'bcrypt hashing (10 salt rounds)\n+ account lockout after 5 failures', RED),
    ('🛡️', 'AES-256 Encryption', 'GCM mode encryption for\nall sensitive data at rest', PURPLE),
    ('⚡', 'Rate Limiting', '100 req/15min general\n5 req/15min auth endpoints', ORANGE),
    ('🪖', 'Helmet Headers', 'CSP, COEP, XSS protection\nsecurity headers middleware', BLUE),
    ('🌐', 'CORS Policy', 'Strict whitelisted origin\npolicy with no wildcards', GREEN),
    ('📋', 'Audit Trail', 'Every API request logged to\nMongoDB + Winston file logs', TEAL),
    ('✅', 'Input Validator', 'express-validator on all\npublic-facing endpoints', CYAN),
    ('🔌', 'Enterprise MW', 'Request IDs, API versioning,\nperformance monitoring', YELLOW),
]

for i, (icon, title, desc, color) in enumerate(security_layers):
    col = i % 5
    row = i // 5
    x = Inches(0.3 + col * 2.55)
    y = Inches(2.0 + row * 2.5)
    c = card(slide, x, y, Inches(2.35), Inches(2.1), BG_CARD)
    accent_line(slide, x, y, Inches(2.35), color)
    txt(slide, x, y + Inches(0.2), Inches(2.35), Inches(0.4), icon, size=22, align=PP_ALIGN.CENTER)
    txt(slide, x, y + Inches(0.6), Inches(2.35), Inches(0.3), title, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.1), y + Inches(1.0), Inches(2.15), Inches(0.9), desc, size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 17: EDUCATION & GAMIFICATION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 15, 'Engagement', 'Education, Gamification & Wellness')

engage_feats = [
    ('📚', 'Learning Center', 'Financial education courses with\nprogressive learning paths and\nquiz-based skill assessment', BLUE),
    ('🧠', 'Financial Quiz', 'Interactive knowledge quizzes\nacross budgeting, investing,\ntax planning, and debt', GREEN),
    ('🏆', 'Achievements', 'Gamified milestone system:\nfirst budget, investment,\n10K saved, consistency streaks', PURPLE),
    ('📍', 'Milestones', 'Net worth milestones, debt\nfreedom progress, portfolio\ngrowth celebrations', ORANGE),
    ('❤️', 'Wellness Score', '8-dimension holistic scoring:\nincome stability, expense mgmt,\nsavings, debt, goals (0-100)', RED),
    ('👥', 'Peer Comparison', 'Anonymous benchmarking by\nincome bracket: "You save\nmore than 72% of peers"', TEAL),
    ('🧪', 'Behavioral Finance', 'Cognitive bias detection:\nloss aversion, anchoring,\nsunk cost, herd behavior', PINK),
    ('🔔', 'Smart Nudges', 'AI-generated behavioral nudges\nto improve financial habits\nbased on spending patterns', CYAN),
]

for i, (icon, title, desc, color) in enumerate(engage_feats):
    col = i % 4
    row = i // 4
    x = Inches(0.8 + col * 3.1)
    y = Inches(2.0 + row * 2.5)
    feature_card(slide, x, y, Inches(2.85), Inches(2.1), icon[0], title, desc, color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 18: TAX, INSURANCE & RETIREMENT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 16, 'Tax & Retirement', 'Tax Optimization, Insurance & Retirement Planning')

# Tax section
card(slide, Inches(0.8), Inches(2.0), Inches(3.8), Inches(4.8), BG_CARD)
accent_line(slide, Inches(0.8), Inches(2.0), Inches(3.8), GREEN)
txt(slide, Inches(1.0), Inches(2.2), Inches(3.4), Inches(0.35), '💰 Tax Features', size=16, color=GREEN, bold=True)

tax_items = [
    'Old vs New Regime Comparison',
    'Section 80C/80D/80E Tracking',
    'AI Tax Optimization Engine',
    'STCG/LTCG Capital Gains Calc',
    'Tax-Loss Harvesting Engine',
    'CII Indexation Benefits',
    'Section 54/54F Exemptions',
    'FY 2025-26 Tax Rules',
]
for i, item in enumerate(tax_items):
    y = Inches(2.7 + i * 0.45)
    txt(slide, Inches(1.2), y, Inches(3.2), Inches(0.3), f'✓  {item}', size=12, color=LIGHT)

# Insurance
card(slide, Inches(4.9), Inches(2.0), Inches(3.8), Inches(4.8), BG_CARD)
accent_line(slide, Inches(4.9), Inches(2.0), Inches(3.8), BLUE)
txt(slide, Inches(5.1), Inches(2.2), Inches(3.4), Inches(0.35), '🛡️ Insurance', size=16, color=BLUE, bold=True)

ins_items = [
    'Insurance Coverage Planner',
    'Life Insurance Calculator',
    'Health Insurance Calculator',
    'Premium Tracking History',
    'Gap Analysis & Recommendations',
    'Policy Portfolio Overview',
]
for i, item in enumerate(ins_items):
    y = Inches(2.7 + i * 0.45)
    txt(slide, Inches(5.3), y, Inches(3.2), Inches(0.3), f'✓  {item}', size=12, color=LIGHT)

# Retirement
card(slide, Inches(9.0), Inches(2.0), Inches(3.5), Inches(4.8), BG_CARD)
accent_line(slide, Inches(9.0), Inches(2.0), Inches(3.5), PURPLE)
txt(slide, Inches(9.2), Inches(2.2), Inches(3.1), Inches(0.35), '🏖️ Retirement', size=16, color=PURPLE, bold=True)

ret_items = [
    'Corpus Planning + Inflation',
    'PPF Tracker (7.1% p.a.)',
    'EPF Tracker',
    'NPS Management',
    'FIRE Calculator',
    'SIP Returns Calculator',
    'Lumpsum Calculator',
    'SIP Delay Cost Analysis',
]
for i, item in enumerate(ret_items):
    y = Inches(2.7 + i * 0.45)
    txt(slide, Inches(9.4), y, Inches(3.0), Inches(0.3), f'✓  {item}', size=12, color=LIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 19: GOALS, NET WORTH & SAVINGS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 17, 'Goals & Wealth', 'Goal Tracking, Net Worth & Savings')

features_grid = [
    ('🎯', 'Financial Goals', 'Create goals with target amount,\ndeadline, priority. Track\ncontributions & milestones.', BLUE),
    ('📈', 'Net Worth Tracker', 'Assets vs liabilities with\nauto-generation from data.\nHistorical trend & projections.', GREEN),
    ('🏠', 'Property Manager', 'Real estate portfolio\nmanagement with valuations\nand rental income tracking.', PURPLE),
    ('💰', 'Savings Challenges', 'Gamified savings with weekly\nand monthly challenges.\nLeaderboards & streaks.', ORANGE),
    ('📊', 'Goal Forecaster', 'AI predicts goal achievement\ndate based on current\ncontribution patterns.', TEAL),
    ('🔥', 'FIRE Tracker', 'Financial Independence Retire\nEarly calculator with\nsafe withdrawal rate.', RED),
    ('💳', 'Credit Score', 'CIBIL score tracking with\ndetailed factor breakdown\nand improvement plan.', PINK),
    ('🏦', 'Wealth Management', 'Comprehensive wealth dashboard\nacross all asset classes\nwith growth projections.', CYAN),
]

for i, (icon, title, desc, color) in enumerate(features_grid):
    col = i % 4
    row = i // 4
    x = Inches(0.8 + col * 3.1)
    y = Inches(2.0 + row * 2.5)
    feature_card(slide, x, y, Inches(2.85), Inches(2.1), icon[0], title, desc, color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 20: THEME SYSTEM & UX
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 18, 'Design', 'Theme System & User Experience')

# Theme modes
modes = [
    ('☀️', 'Light', 'Default bright theme\nwith clean typography', RGBColor(0xF8, 0xFA, 0xFC)),
    ('🌙', 'Dark', 'Reduced eye strain\nfor night usage', RGBColor(0x1E, 0x29, 0x3B)),
    ('⬛', 'AMOLED', 'Pure black for OLED\nscreens — saves battery', RGBColor(0x00, 0x00, 0x00)),
]

for i, (icon, name, desc, preview_color) in enumerate(modes):
    x = Inches(0.8 + i * 2.2)
    # Preview card
    c = card(slide, x, Inches(2.0), Inches(2.0), Inches(1.8), preview_color)
    if i > 0:
        txt(slide, x, Inches(2.2), Inches(2.0), Inches(0.4), icon, size=24, color=WHITE, align=PP_ALIGN.CENTER)
    else:
        txt(slide, x, Inches(2.2), Inches(2.0), Inches(0.4), icon, size=24, color=BG_DARK, align=PP_ALIGN.CENTER)
    txt(slide, x, Inches(2.7), Inches(2.0), Inches(0.3), name, size=14, color=WHITE if i > 0 else BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.1), Inches(3.05), Inches(1.8), Inches(0.6), desc, size=10, color=MUTED if i > 0 else DIM, align=PP_ALIGN.CENTER)

# Accent colors
card(slide, Inches(7.2), Inches(2.0), Inches(5.3), Inches(1.8), BG_CARD)
accent_line(slide, Inches(7.2), Inches(2.0), Inches(5.3), BLUE)
txt(slide, Inches(7.4), Inches(2.1), Inches(5.0), Inches(0.35), '8 Accent Color Presets', size=14, color=WHITE, bold=True)

accents = [
    ('Ocean', BLUE), ('Purple', PURPLE), ('Emerald', GREEN), ('Rose', PINK),
    ('Amber', ORANGE), ('Teal', TEAL), ('Indigo', RGBColor(0x63, 0x66, 0xF1)), ('Sky', CYAN),
]
for i, (name, color) in enumerate(accents):
    x = Inches(7.5 + i * 0.6)
    circle(slide, x, Inches(2.6), Inches(0.4), color)
    txt(slide, x - Inches(0.1), Inches(3.1), Inches(0.6), Inches(0.25), name, size=8, color=MUTED, align=PP_ALIGN.CENTER)

# UX features
ux_feats = [
    ('⌨️', 'Keyboard Shortcuts', 'Ctrl+K search, Ctrl+N new txn,\nglobal shortcuts with guide', BLUE),
    ('📱', 'Responsive Design', 'Mobile-first responsive with\ndedicated mobile CSS layers', GREEN),
    ('🤖', 'Smart Assistant', 'Floating AI assistant panel\navailable on every page', PURPLE),
    ('📌', 'Sidebar Navigation', '9 collapsible sections, auto-expand\non route, localStorage persistence', ORANGE),
    ('🔔', 'Notification Bell', 'Real-time notification count\nwith WebSocket updates', RED),
    ('⚡', 'Quick Add', 'Navbar expense button + floating\nFAB for instant transaction entry', TEAL),
]

for i, (icon, title, desc, color) in enumerate(ux_feats):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(4.3 + row * 1.5)
    feature_card(slide, x, y, Inches(3.8), Inches(1.2), icon[0], title, desc, color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 21: KEY DIFFERENTIATORS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 19, 'Differentiators', 'What Makes Us Different')

# Big differentiator cards
diffs = [
    ('🔒', '100% Local AI', 'Every ML algorithm runs locally on your server.\nZero data leaves your infrastructure.\nNo OpenAI/Claude API costs — unlimited predictions.',
     PURPLE, '33,000+ lines of from-scratch ML code'),
    ('🇮🇳', 'Indian Market', 'INR formatting (Cr/L/K), CIBIL scoring, Indian\ntax rules FY 25-26, UPI detection, NPS/PPF/ELSS\nasset classes, 55+ Indian subscription services.',
     GREEN, 'Built for India from day one'),
    ('📊', '3-Gen UI', 'V1 (Functional) → V2 (Enhanced MUI) → V3 (Enterprise).\nUsers pick their complexity level. Progressive\nenhancement without breaking existing workflows.',
     BLUE, 'Choose your experience level'),
    ('🧠', 'Self-Learning', 'Per-user models retrain automatically. Drift\ndetection triggers re-learning. A/B testing for\nmodel versions. Zero manual ML operations.',
     ORANGE, 'Gets smarter with every transaction'),
]

for i, (icon, title, desc, color, tagline) in enumerate(diffs):
    y = Inches(2.0 + i * 1.3)
    c = card(slide, Inches(0.8), y, Inches(11.7), Inches(1.15), BG_CARD)
    accent_line(slide, Inches(0.8), y, Inches(0.07), color, Inches(1.15))
    icon_circle(slide, Inches(1.1), y + Inches(0.3), icon[0], color, Inches(0.55))
    txt(slide, Inches(1.9), y + Inches(0.1), Inches(3.5), Inches(0.35), title, size=18, color=WHITE, bold=True)
    txt(slide, Inches(1.9), y + Inches(0.45), Inches(6.0), Inches(0.7), desc, size=11, color=MUTED)
    # Right tagline
    pill_card = card(slide, Inches(9.0), y + Inches(0.35), Inches(3.3), Inches(0.35), color)
    txt(slide, Inches(9.0), y + Inches(0.35), Inches(3.3), Inches(0.35), tagline, size=10, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 22: SUMMARY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
gradient_bar(slide, Inches(0), Inches(0), W, Inches(0.08), BLUE, PURPLE)

txt(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
    'EXECUTIVE SUMMARY', size=14, color=BLUE, bold=True)
txt(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
    'Platform at a Glance', size=36, color=WHITE, bold=True)

# Summary stats in a bigger grid
summary_stats = [
    ('640+', 'Source Files', BLUE),
    ('150+', 'Features', GREEN),
    ('300+', 'API Endpoints', PURPLE),
    ('33K+', 'AI Code Lines', ORANGE),
    ('38', 'ML Modules', TEAL),
    ('25+', 'Algorithms', PINK),
    ('146', 'Frontend Pages', CYAN),
    ('47', 'DB Models', YELLOW),
    ('0', 'External API\nDependencies', RED),
]

for i, (num, label, color) in enumerate(summary_stats):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.8 + row * 1.7)
    stat_box(slide, x, y, Inches(3.8), Inches(1.4), num, label, color)

# Bottom key message
card(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.7), BG_SURFACE)
txt(slide, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.6),
    'The most comprehensive AI-powered financial management platform — 100% local, Indian market focused, enterprise-ready.',
    size=15, color=LIGHT, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 23: THANK YOU
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
gradient_bar(slide, Inches(0), Inches(0), W, Inches(0.08), BLUE, PURPLE)

# Decorative circles (faint)
for cx, cy, sz, clr in [(10.5, 1.0, 3.0, RGBColor(0x10, 0x16, 0x25)), (9.5, 4.0, 3.5, RGBColor(0x12, 0x14, 0x28)), (1.0, 5.0, 2.5, RGBColor(0x0E, 0x16, 0x20))]:
    circle(slide, Inches(cx), Inches(cy), Inches(sz), clr)

txt(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.2),
    'Thank You', size=60, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Accent divider
gradient_bar(slide, Inches(5.0), Inches(3.3), Inches(3.3), Inches(0.05), BLUE, PURPLE)

txt(slide, Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.6),
    'FinancialAnalyzer v2.0.0', size=24, color=MUTED, align=PP_ALIGN.CENTER)

txt(slide, Inches(0.8), Inches(4.4), Inches(11.7), Inches(0.4),
    'by Circuvent Technologies', size=16, color=BLUE, align=PP_ALIGN.CENTER)

# Contact info card
card(slide, Inches(3.5), Inches(5.2), Inches(6.3), Inches(1.2), BG_SURFACE)
txt(slide, Inches(3.5), Inches(5.3), Inches(6.3), Inches(0.35),
    'Enterprise Financial Management  •  100% Local AI  •  Indian Market Focus', size=13, color=MUTED, align=PP_ALIGN.CENTER)
txt(slide, Inches(3.5), Inches(5.7), Inches(6.3), Inches(0.35),
    'support@circuvent.com  •  www.circuvent.com', size=12, color=DIM, align=PP_ALIGN.CENTER)
txt(slide, Inches(3.5), Inches(6.05), Inches(6.3), Inches(0.3),
    'Questions & Demo Request', size=14, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# Footer
txt(slide, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.3),
    '© 2026 Circuvent Technologies. All rights reserved.', size=10, color=FAINT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output = r'c:\Users\v-hbonthada\WorkSpace-Pract\FinancialAnalyzer\docs\FinancialAnalyzer-Business-Presentation.pptx'
prs.save(output)
print(f'✅ Presentation saved: {output}')
print(f'📊 Total slides: {len(prs.slides)}')
