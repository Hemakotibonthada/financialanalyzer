"""
FinancialAnalyzer — Comprehensive Feature Presentation Generator
Generates a professional PowerPoint presentation covering all features.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── COLOR PALETTE ──
NAVY      = RGBColor(0x0F, 0x17, 0x2A)
DARK_BLUE = RGBColor(0x1E, 0x3A, 0x5F)
BLUE      = RGBColor(0x2D, 0x6A, 0x9F)
ACCENT    = RGBColor(0x38, 0xBD, 0xF8)
GREEN     = RGBColor(0x10, 0xB9, 0x81)
PURPLE    = RGBColor(0x8B, 0x5C, 0xF6)
ORANGE    = RGBColor(0xF5, 0x9E, 0x0B)
RED       = RGBColor(0xEF, 0x44, 0x44)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0xE2, 0xE8, 0xF0)
GRAY      = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY = RGBColor(0x47, 0x55, 0x69)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=ACCENT):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        p.space_after = Pt(6)
        p.level = 0
    return tf


def add_stat_card(slide, left, top, width, height, number, label, color=ACCENT):
    card = add_shape(slide, left, top, width, height, DARK_BLUE)
    card.fill.fore_color.rgb = DARK_BLUE
    # Number
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.8),
                 str(number), font_size=32, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, left + Inches(0.2), top + Inches(0.9), width - Inches(0.4), Inches(0.5),
                 label, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)
    return card


def add_feature_row(slide, left, top, width, icon_text, title, description, color=ACCENT):
    # Icon circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Inches(0.05), Inches(0.4), Inches(0.4))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].text = icon_text
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.bold = True
    # Title
    add_text_box(slide, left + Inches(0.55), top, width - Inches(0.6), Inches(0.3),
                 title, font_size=16, color=WHITE, bold=True)
    # Description
    add_text_box(slide, left + Inches(0.55), top + Inches(0.28), width - Inches(0.6), Inches(0.3),
                 description, font_size=12, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide)

# Accent bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

# Company
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
             'CIRCUVENT TECHNOLOGIES', font_size=16, color=ACCENT, bold=True)

# Title
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1.5),
             'FinancialAnalyzer', font_size=54, color=WHITE, bold=True)

# Subtitle
add_text_box(slide, Inches(0.8), Inches(3.2), Inches(11), Inches(0.8),
             'Enterprise Financial Management Platform with AI-First Architecture',
             font_size=24, color=GRAY)

# Divider
add_shape(slide, Inches(0.8), Inches(4.2), Inches(3), Inches(0.04), ACCENT)

# Stats row
add_stat_card(slide, Inches(0.8), Inches(4.8), Inches(2.2), Inches(1.3), '640+', 'Source Files')
add_stat_card(slide, Inches(3.3), Inches(4.8), Inches(2.2), Inches(1.3), '150+', 'Features', GREEN)
add_stat_card(slide, Inches(5.8), Inches(4.8), Inches(2.2), Inches(1.3), '300+', 'API Endpoints', PURPLE)
add_stat_card(slide, Inches(8.3), Inches(4.8), Inches(2.2), Inches(1.3), '33K+', 'Lines AI Code', ORANGE)

# Version
add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
             'Version 2.0.0  •  March 2026  •  100% Local AI  •  Zero External API Dependencies',
             font_size=14, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: TABLE OF CONTENTS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             'PRESENTATION OVERVIEW', font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             'Table of Contents', font_size=36, color=WHITE, bold=True)

toc_items = [
    ('01', 'Technology Stack & Architecture'),
    ('02', 'Core Financial Features'),
    ('03', 'Transaction Management'),
    ('04', 'Budget & Planning'),
    ('05', 'Investment & Portfolio'),
    ('06', 'Debt, EMI & Loan Management'),
    ('07', 'AI & Machine Learning Engine'),
    ('08', 'AI Lab — Local ML Features'),
    ('09', 'NLP & Conversational AI'),
    ('10', 'Anomaly Detection & Fraud Prevention'),
    ('11', 'Enterprise V2/V3 Features'),
    ('12', 'Reports, Analytics & Export'),
    ('13', 'Integrations & Platform'),
    ('14', 'Security & Authentication'),
    ('15', 'Summary & Key Differentiators'),
]

col1 = toc_items[:8]
col2 = toc_items[8:]

for i, (num, title) in enumerate(col1):
    y = Inches(1.8) + Inches(i * 0.55)
    add_text_box(slide, Inches(1.0), y, Inches(0.5), Inches(0.4), num, font_size=18, color=ACCENT, bold=True)
    add_text_box(slide, Inches(1.6), y, Inches(4.5), Inches(0.4), title, font_size=16, color=LIGHT_GRAY)

for i, (num, title) in enumerate(col2):
    y = Inches(1.8) + Inches(i * 0.55)
    add_text_box(slide, Inches(7.0), y, Inches(0.5), Inches(0.4), num, font_size=18, color=ACCENT, bold=True)
    add_text_box(slide, Inches(7.6), y, Inches(4.5), Inches(0.4), title, font_size=16, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: TECH STACK
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
             '01  TECHNOLOGY STACK', font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             'Full-Stack Architecture', font_size=36, color=WHITE, bold=True)

# Backend card
card = add_shape(slide, Inches(0.8), Inches(1.8), Inches(3.7), Inches(4.8), DARK_BLUE)
add_text_box(slide, Inches(1.0), Inches(1.9), Inches(3.3), Inches(0.4), '⚙️  Backend', font_size=20, color=GREEN, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(2.4), Inches(3.3), Inches(4.0), [
    '• Node.js 18+ Runtime',
    '• Express.js 4.x Framework',
    '• MongoDB 6+ (Mongoose ODM)',
    '• Socket.IO 4.x Real-time',
    '• Redis Cache (optional)',
    '• JWT + bcrypt Auth',
    '• Speakeasy 2FA (TOTP)',
    '• Winston Logging',
    '• 89 Route Files',
    '• 128 Service Files',
    '• 47 Mongoose Models',
], font_size=13)

# Frontend card
card = add_shape(slide, Inches(4.8), Inches(1.8), Inches(3.7), Inches(4.8), DARK_BLUE)
add_text_box(slide, Inches(5.0), Inches(1.9), Inches(3.3), Inches(0.4), '🎨  Frontend', font_size=20, color=ACCENT, bold=True)
add_bullet_list(slide, Inches(5.0), Inches(2.4), Inches(3.3), Inches(4.0), [
    '• React 18 + Vite 5',
    '• Material UI 7',
    '• Tailwind CSS 3',
    '• Chart.js 4 + Recharts 2',
    '• Lucide React Icons',
    '• React Router 6',
    '• Socket.IO Client',
    '• Firebase 12 (Auth)',
    '• 146 Page Components',
    '• 67 UI Components',
    '• 8 React Contexts',
], font_size=13)

# AI card
card = add_shape(slide, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8), DARK_BLUE)
add_text_box(slide, Inches(9.0), Inches(1.9), Inches(3.3), Inches(0.4), '🧠  AI/ML (100% Local)', font_size=20, color=PURPLE, bold=True)
add_bullet_list(slide, Inches(9.0), Inches(2.4), Inches(3.3), Inches(4.0), [
    '• Neural Networks (from scratch)',
    '• Random Forest / GBT',
    '• ARIMA / Holt-Winters',
    '• K-Means++ / DBSCAN / PCA',
    '• TF-IDF / NER / Sentiment',
    '• Q-Learning / DQN / A-C',
    '• Isolation Forest / LOF',
    '• SHAP / LIME / XAI',
    '• Markowitz Optimization',
    '• 38 AI Modules',
    '• ~33,000 Lines of Code',
], font_size=13)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
             '01  ARCHITECTURE', font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             'System Architecture Overview', font_size=36, color=WHITE, bold=True)

# Architecture layers
layers = [
    ('Client Layer', 'React SPA • Electron Desktop • Mobile (React Native)', ACCENT, Inches(1.8)),
    ('API Gateway', 'Express.js • Helmet • Rate Limiting • CORS • JWT Auth • Enterprise MW', GREEN, Inches(2.6)),
    ('Route Layer', '89 Route Files • 70+ Prefixes • 300+ Endpoints', BLUE, Inches(3.4)),
    ('Service Layer', '82 Core Services + 46 AI Modules = 128 Services', PURPLE, Inches(4.2)),
    ('AI/ML Engine', '38 AI Modules • Neural Networks • NLP • RL • Anomaly Detection', ORANGE, Inches(5.0)),
    ('Data Layer', 'MongoDB (47 Models) • Redis Cache • File System (ML Models)', RED, Inches(5.8)),
]

for label, desc, color, top in layers:
    bar = add_shape(slide, Inches(1.0), top, Inches(0.1), Inches(0.6), color)
    add_text_box(slide, Inches(1.3), top, Inches(3.0), Inches(0.35), label, font_size=18, color=color, bold=True)
    add_text_box(slide, Inches(1.3), top + Inches(0.3), Inches(10.0), Inches(0.3), desc, font_size=13, color=GRAY)

# Right side features
add_shape(slide, Inches(9.0), Inches(1.8), Inches(3.5), Inches(4.8), DARK_BLUE)
add_text_box(slide, Inches(9.2), Inches(1.9), Inches(3.0), Inches(0.4), 'Key Patterns', font_size=18, color=ACCENT, bold=True)
add_bullet_list(slide, Inches(9.2), Inches(2.4), Inches(3.0), Inches(4.0), [
    '✓ 100% Local AI',
    '✓ Zero API Dependencies',
    '✓ Per-User ML Models',
    '✓ Self-Learning Pipeline',
    '✓ Drift Detection',
    '✓ Dual Auth Backend',
    '✓ Multi-Version UI (V1/V2/V3)',
    '✓ Indian Market Focus',
    '✓ Real-time WebSocket',
    '✓ PWA + Desktop + Mobile',
], font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: CORE FINANCIAL FEATURES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GREEN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
             '02  CORE FEATURES', font_size=14, color=GREEN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             'Core Financial Features', font_size=36, color=WHITE, bold=True)

features = [
    ('💰', 'Dashboard', '3 versions (V1/V2/V3) with KPIs, spending trends, AI predictions', GREEN),
    ('📊', 'Financial Health', 'Health score 0-100, breakdown across savings, debt, insurance', GREEN),
    ('🏦', 'Bank Accounts', 'Multi-bank management with balance tracking', BLUE),
    ('💳', 'Credit Cards', 'Card tracking, bill management, statement parsing', BLUE),
    ('💱', 'Currency', '8 currencies with live conversion rates', ACCENT),
    ('👨‍👩‍👧‍👦', 'Family Finance', 'Family member management and shared budgets', PURPLE),
    ('🏢', 'Business Suite', 'Invoices, clients, contracts, projects, vendors', ORANGE),
    ('🎮', 'Gamification', 'Achievements, milestones, savings challenges', ACCENT),
]

for i, (icon, title, desc, color) in enumerate(features):
    row = i % 4
    col = i // 4
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(1.8) + Inches(row * 1.2)
    add_feature_row(slide, x, y, Inches(5.8), icon[0] if len(icon) == 1 else '★', title, desc, color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: TRANSACTIONS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GREEN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
             '03  TRANSACTIONS', font_size=14, color=GREEN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             'Transaction Management', font_size=36, color=WHITE, bold=True)

# Left: features
add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), DARK_BLUE)
add_text_box(slide, Inches(1.0), Inches(1.9), Inches(5.0), Inches(0.4), 'Input Sources', font_size=20, color=ACCENT, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(2.4), Inches(5.0), Inches(4.0), [
    '• Manual Entry — Quick add with category, amount, merchant',
    '• Document Upload — PDF/CSV/image bank statement parsing',
    '• Gmail Sync — Auto-extract from bank email alerts',
    '• Receipt Scanner — OCR-powered receipt scanning',
    '• CSV Import — Bulk import with column mapping',
    '• UPI Detection — Extract VPA, UTR, app details',
    '',
    'Payment Methods:',
    '  Cash • Card • UPI • Bank Transfer • Wallet',
    '  Net Banking • Cheque • IMPS • NEFT • RTGS',
], font_size=14)

# Right: AI pipeline
add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0), DARK_BLUE)
add_text_box(slide, Inches(7.0), Inches(1.9), Inches(5.0), Inches(0.4), 'AI Enrichment Pipeline', font_size=20, color=PURPLE, bold=True)
add_bullet_list(slide, Inches(7.0), Inches(2.4), Inches(5.0), Inches(4.0), [
    '1. Category Classification (ML model)',
    '2. Merchant Name Normalization',
    '3. Anomaly Score (Isolation Forest)',
    '4. UPI Metadata Extraction',
    '5. Recurring Pattern Detection',
    '6. Tax Deductibility Check',
    '7. Budget Impact Calculation',
    '8. Notification Triggers',
    '',
    'Full-text search with weighted indexing',
    'Advanced filters: date, category, type, method',
], font_size=14)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: BUDGET & PLANNING
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GREEN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
             '04  BUDGET & PLANNING', font_size=14, color=GREEN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             'Budget Management & Financial Planning', font_size=36, color=WHITE, bold=True)

budget_features = [
    ('📋', 'Category Budgets', 'Set monthly/weekly/yearly budgets per expense category', GREEN),
    ('📊', 'Auto-Spending Calc', 'Real-time spent amount calculated from transactions', BLUE),
    ('⚠️', 'Alert Thresholds', 'Configurable alerts at 80% (default) budget usage', ORANGE),
    ('🔄', 'Budget Rollover', 'Carry forward unused budget to next period', ACCENT),
    ('✨', '50/30/20 Auto-Fill', 'Auto-calculate budgets from salary (needs/wants/savings)', PURPLE),
    ('🧙', 'Smart Budget Wizard', 'Step-by-step guided budget creation', ACCENT),
    ('🤖', 'AI Budget Optimizer', 'Reinforcement Learning-powered optimal allocation', PURPLE),
    ('📅', 'Financial Calendar', 'Calendar view of all upcoming financial events', BLUE),
]

for i, (icon, title, desc, color) in enumerate(budget_features):
    row = i % 4
    col = i // 4
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(1.8) + Inches(row * 1.2)
    add_feature_row(slide, x, y, Inches(5.8), '★', title, desc, color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: INVESTMENTS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
             '05  INVESTMENTS', font_size=14, color=BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             'Investment Portfolio & Market Features', font_size=36, color=WHITE, bold=True)

# Left
add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), DARK_BLUE)
add_text_box(slide, Inches(1.0), Inches(1.9), Inches(5.0), Inches(0.4), 'Investment Types', font_size=20, color=ACCENT, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(2.4), Inches(5.0), Inches(4.2), [
    '• Stocks & Equities',
    '• Mutual Funds (with NAV tracking)',
    '• Fixed Deposits (with maturity alerts)',
    '• Gold (physical + digital)',
    '• Crypto Portfolio',
    '• PPF / EPF / NPS',
    '• ELSS (Tax Saving)',
    '• REITs (Real Estate)',
    '• Bonds & Debentures',
    '• SIP (Systematic Investment Plans)',
], font_size=14)

# Right: Analytics
add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0), DARK_BLUE)
add_text_box(slide, Inches(7.0), Inches(1.9), Inches(5.0), Inches(0.4), 'Analytics (22 Endpoints)', font_size=20, color=GREEN, bold=True)
add_bullet_list(slide, Inches(7.0), Inches(2.4), Inches(5.0), Inches(4.2), [
    '• Portfolio Dashboard & Summary',
    '• XIRR, CAGR, Absolute Returns',
    '• Asset Allocation Analysis',
    '• Risk Analysis & VaR',
    '• Diversification Score (0-100)',
    '• Tax Efficiency Analysis',
    '• AI Recommendations (Buy/Sell/Hold)',
    '• Markowitz Portfolio Optimization',
    '• Efficient Frontier Generation',
    '• Monte Carlo Simulation',
], font_size=14)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: DEBT & EMI
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), RED)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
             '06  DEBT & EMI', font_size=14, color=RED, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             'Debt, EMI & Loan Management', font_size=36, color=WHITE, bold=True)

# Stats
add_stat_card(slide, Inches(0.8), Inches(1.8), Inches(2.5), Inches(1.2), '33', 'EMI Endpoints', RED)
add_stat_card(slide, Inches(3.6), Inches(1.8), Inches(2.5), Inches(1.2), '3', 'UI Versions', ORANGE)
add_stat_card(slide, Inches(6.4), Inches(1.8), Inches(2.5), Inches(1.2), '4384', 'Lines (EMI Routes)', PURPLE)
add_stat_card(slide, Inches(9.2), Inches(1.8), Inches(2.5), Inches(1.2), '8', 'Debt Types', ACCENT)

# Features grid
debt_features = [
    ('EMI Tracker', 'Monthly tracking with principal/interest split, statement sync'),
    ('Debt Payoff', 'Snowball vs Avalanche vs AI Hybrid comparison'),
    ('Foreclosure', 'Calculate costs & savings of early closure'),
    ('One-Click Prepay', 'Prepayment intent with auto surplus sweep'),
    ('Balance Transfer', 'Analyze transfer opportunities for lower rates'),
    ('Late Fee Shield', 'Automated late fee protection system'),
    ('Emergency Fund', 'Track emergency fund with contribution history'),
    ('Lender Dashboard', 'Role-based dashboard for money lenders'),
    ('Loan Calculator', 'EMI/amortization calculation with schedule'),
    ('Multi-Currency', 'EMI tracking in USD/INR with exchange rates'),
    ('PDF/Excel Export', 'Export EMI schedules in multiple formats'),
    ('Bank Deduction', 'Track auto-deduction from bank accounts'),
]

for i, (title, desc) in enumerate(debt_features):
    row = i % 4
    col = i // 4
    x = Inches(0.8) + Inches(col * 4.2)
    y = Inches(3.4) + Inches(row * 0.95)
    add_text_box(slide, x, y, Inches(3.8), Inches(0.3), f'▸ {title}', font_size=14, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.28), Inches(3.6), Inches(0.4), desc, font_size=11, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: AI/ML ENGINE OVERVIEW
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
             '07  AI/ML ENGINE', font_size=14, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             'AI & Machine Learning — 100% Local', font_size=36, color=WHITE, bold=True)

# Stats
add_stat_card(slide, Inches(0.8), Inches(1.7), Inches(2.0), Inches(1.1), '38', 'AI Modules', PURPLE)
add_stat_card(slide, Inches(3.1), Inches(1.7), Inches(2.0), Inches(1.1), '33K+', 'Lines of Code', ACCENT)
add_stat_card(slide, Inches(5.4), Inches(1.7), Inches(2.0), Inches(1.1), '100+', 'Endpoints', GREEN)
add_stat_card(slide, Inches(7.7), Inches(1.7), Inches(2.0), Inches(1.1), '25+', 'Algorithms', ORANGE)
add_stat_card(slide, Inches(10.0), Inches(1.7), Inches(2.0), Inches(1.1), '0', 'External APIs', RED)

# Categories
ai_cats = [
    ('Core ML', 'Neural Networks, Decision Trees, Clustering, Time Series', '4,650 lines', ACCENT),
    ('Training', 'Pipeline, Self-Learning, AutoML, Orchestrator, Monitoring', '4,880 lines', GREEN),
    ('Prediction', 'Financial Forecasting, Credit Score, Cash Flow', '2,360 lines', BLUE),
    ('Anomaly', 'Isolation Forest, LOF, Autoencoders, Fraud Detection', '2,510 lines', RED),
    ('NLP', 'TF-IDF, NER, Sentiment, Chatbot, Semantic Search', '4,270 lines', PURPLE),
    ('Optimization', 'RL (DQN/Actor-Critic), Portfolio, Tax, Goals', '3,970 lines', ORANGE),
    ('Behavioral', 'Bias Detection, Spending Intel, Peer Compare, Wellness', '3,200 lines', GREEN),
    ('XAI', 'SHAP, LIME, Counterfactual, Knowledge Graph', '1,805 lines', ACCENT),
]

for i, (cat, desc, lines, color) in enumerate(ai_cats):
    row = i % 4
    col = i // 4
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(3.2) + Inches(row * 1.0)
    bar = add_shape(slide, x, y + Inches(0.05), Inches(0.08), Inches(0.5), color)
    add_text_box(slide, x + Inches(0.2), y, Inches(2.5), Inches(0.3), cat, font_size=15, color=color, bold=True)
    add_text_box(slide, x + Inches(2.7), y, Inches(1.0), Inches(0.3), lines, font_size=12, color=GRAY, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(5.5), Inches(0.35), desc, font_size=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11: AI LAB
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
             '08  AI LAB', font_size=14, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             '13 Dedicated AI Feature Pages', font_size=36, color=WHITE, bold=True)

ai_lab = [
    ('AI Chatbot', 'Multi-turn conversational financial AI assistant'),
    ('RL Optimizer', 'Budget/investment optimization using Q-Learning, DQN'),
    ('Model Observatory', 'ML model performance monitoring, drift detection'),
    ('Anomaly Lab', 'Multi-algorithm anomaly detection (IF, LOF, SPC)'),
    ('Financial Planner', 'AI-powered comprehensive financial planning'),
    ('Spending Intel', 'Merchant-level insights, impulse detection'),
    ('Portfolio AI', 'Markowitz, Black-Litterman, Efficient Frontier'),
    ('Credit Predictor', 'CIBIL score prediction and what-if simulation'),
    ('Cash Flow AI', 'Cash flow analysis, income patterns, liquidity'),
    ('Sub Manager', 'Auto-detect 55+ Indian subscriptions, optimization'),
    ('Goal & Tax AI', 'Goal achievement + tax harvesting (STCG/LTCG)'),
    ('Wellness AI', '8-dimension wellness scoring'),
    ('AI Command V3', 'Central AI control panel with health monitoring'),
]

for i, (title, desc) in enumerate(ai_lab):
    if i < 7:
        x = Inches(0.8)
        y = Inches(1.8) + Inches(i * 0.72)
    else:
        x = Inches(7.0)
        y = Inches(1.8) + Inches((i - 7) * 0.72)
    
    add_text_box(slide, x, y, Inches(5.5), Inches(0.3), f'🧠  {title}', font_size=15, color=ACCENT, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(5.0), Inches(0.3), desc, font_size=12, color=GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12: NLP & CONVERSATIONAL AI
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
             '09  NLP & CONVERSATIONAL', font_size=14, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             'Natural Language Processing & AI Chat', font_size=36, color=WHITE, bold=True)

# NLP features
add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), DARK_BLUE)
add_text_box(slide, Inches(1.0), Inches(1.9), Inches(5.0), Inches(0.4), 'NLP Capabilities', font_size=20, color=ACCENT, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(2.4), Inches(5.0), Inches(4.2), [
    '• Tokenization with financial term preservation',
    '• TF-IDF Vectorization for document similarity',
    '• Sentiment Analysis on transactions/descriptions',
    '• Financial NER (amounts ₹, dates, accounts)',
    '• Query Understanding ("show food > ₹500")',
    '• Text Summarization for financial reports',
    '• Semantic Search with fuzzy matching',
    '• Natural Language Report Generator',
    '• Document Intelligence (bank statements, salary slips)',
], font_size=14)

# Conversational AI
add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0), DARK_BLUE)
add_text_box(slide, Inches(7.0), Inches(1.9), Inches(5.0), Inches(0.4), 'Conversational AI', font_size=20, color=GREEN, bold=True)
add_bullet_list(slide, Inches(7.0), Inches(2.4), Inches(5.0), Inches(4.2), [
    '• Multi-turn dialogue with memory',
    '• Short-term + long-term conversation memory',
    '• Intent classification (40+ financial intents)',
    '• Entity extraction (amounts, dates, categories)',
    '• Context tracking across turns',
    '• Preference learning from interactions',
    '• 3 chat versions (V1, V2, V3)',
    '• Example: "How much did I spend on Swiggy?"',
    '• Example: "Compare my spending this vs last month"',
], font_size=14)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13: ANOMALY & FRAUD
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), RED)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
             '10  ANOMALY & FRAUD', font_size=14, color=RED, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             'Anomaly Detection & Fraud Prevention', font_size=36, color=WHITE, bold=True)

# Anomaly
add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), DARK_BLUE)
add_text_box(slide, Inches(1.0), Inches(1.9), Inches(5.0), Inches(0.4), 'Anomaly Detection (5 Algorithms)', font_size=18, color=RED, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(2.4), Inches(5.0), Inches(4.2), [
    '1. Isolation Forest',
    '   Random partitioning — anomalies isolate faster',
    '',
    '2. Local Outlier Factor (LOF)',
    '   Local density comparison — finds local outliers',
    '',
    '3. Statistical Process Control',
    '   Control charts, Shewhart rules',
    '',
    '4. Autoencoder',
    '   Reconstruction error scoring',
    '',
    '5. DBSCAN Clustering',
    '   Density-based outlier detection',
], font_size=13)

# Fraud
add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0), DARK_BLUE)
add_text_box(slide, Inches(7.0), Inches(1.9), Inches(5.0), Inches(0.4), 'Fraud Prevention (5 Layers)', font_size=18, color=ORANGE, bold=True)
add_bullet_list(slide, Inches(7.0), Inches(2.4), Inches(5.0), Inches(4.2), [
    '1. Rule Engine',
    '   Configurable rules (high-value > 5x avg)',
    '',
    '2. ML Anomaly Scoring',
    '   Integration with anomaly detection ensemble',
    '',
    '3. Behavioral Biometrics',
    '   Spending pattern deviation analysis',
    '',
    '4. Velocity Analysis',
    '   Transaction frequency + daily limit checks',
    '',
    '5. Geolocation Analysis',
    '   Location-based anomaly detection',
], font_size=13)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14: ENTERPRISE V2/V3
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
             '11  ENTERPRISE', font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             'Enterprise V2 & V3 Features', font_size=36, color=WHITE, bold=True)

# V2
add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.2), DARK_BLUE)
add_text_box(slide, Inches(1.0), Inches(1.9), Inches(5.0), Inches(0.4), 'Enhanced V2 (11 Pages)', font_size=18, color=ACCENT, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(2.3), Inches(5.0), Inches(1.5), [
    'Dashboard V2 • Health V2 • Spending V2 • Portfolio V2',
    'Budget V2 • Debt V2 • Goals V2 • Transactions V2',
    'AI Chat V2 • Planning V2 • Reports V2 • Settings V2',
], font_size=13)

# V3
add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.2), DARK_BLUE)
add_text_box(slide, Inches(7.0), Inches(1.9), Inches(5.0), Inches(0.4), 'Enterprise V3 (13 Pages)', font_size=18, color=GREEN, bold=True)
add_bullet_list(slide, Inches(7.0), Inches(2.3), Inches(5.0), Inches(1.5), [
    'Dashboard V3 • Transactions V3 • Budget Intelligence',
    'Investment Advisor • Goals V3 • Debt V3 • Health V3',
    'Analytics V3 • Reports V3 • Settings V3 • AI Chat V3',
    'Cashflow Forecaster • Gmail Browser',
], font_size=13)

# Enterprise features
enterprise_list = [
    ('Admin Panel', 'User management, system monitoring, role control', RED),
    ('Lender Dashboard', 'Loan portfolio, borrower management, payments', ORANGE),
    ('Automation Engine', 'If-then rules, triggers, scheduled actions', PURPLE),
    ('Smart Notifications', 'AI priority engine (P0-P3), fatigue prevention', ACCENT),
    ('Enterprise Reports', 'Template-based, scheduled PDF/Excel/CSV delivery', GREEN),
    ('Activity Logs', 'Full audit trail with request/response metadata', BLUE),
]

for i, (title, desc, color) in enumerate(enterprise_list):
    row = i % 3
    col = i // 3
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(4.4) + Inches(row * 0.9)
    add_feature_row(slide, x, y, Inches(5.8), '★', title, desc, color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 15: REPORTS & ANALYTICS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
             '12  REPORTS & ANALYTICS', font_size=14, color=BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             'Reports, Analytics & Data Export', font_size=36, color=WHITE, bold=True)

report_features = [
    ('📊', 'Reports Hub', '3 versions of financial reports with drill-down', BLUE),
    ('📈', 'Advanced Analytics', 'Multi-dimensional analysis with trend spotting', GREEN),
    ('🔬', 'Data Viz Lab', 'Interactive data visualization playground', PURPLE),
    ('📋', 'Export Center', 'PDF, Excel, CSV export with templates', ACCENT),
    ('⚖️', 'Comparison Tool', 'Period-over-period comparison (MoM, YoY)', ORANGE),
    ('🏆', 'Financial Scorecard', 'Comprehensive scoring across all metrics', GREEN),
    ('📉', 'Risk Dashboard', 'Financial risk assessment dashboard', RED),
    ('🔍', 'Spending Insights', 'Deep spending analysis (2 versions)', BLUE),
]

for i, (icon, title, desc, color) in enumerate(report_features):
    row = i % 4
    col = i // 4
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(1.8) + Inches(row * 1.2)
    add_feature_row(slide, x, y, Inches(5.8), '★', title, desc, color)

# Export formats
add_shape(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8), DARK_BLUE)
add_text_box(slide, Inches(1.0), Inches(6.1), Inches(11.0), Inches(0.5),
             'Export Formats:   PDF (with charts)   •   Excel (.xlsx multi-sheet)   •   CSV (raw data)   •   JSON (API-compatible)',
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 16: INTEGRATIONS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
             '13  INTEGRATIONS', font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             'Integrations & Platform Capabilities', font_size=36, color=WHITE, bold=True)

integrations = [
    ('📧', 'Gmail API', 'Auto-sync transactions from bank emails (UPI, NEFT, CC alerts)', GREEN),
    ('☁️', 'Google Drive', 'Cloud backup and restore of financial data', BLUE),
    ('🏦', 'Plaid Banking', 'Direct bank account connection and transaction import', ACCENT),
    ('📱', 'Firebase', 'Google OAuth authentication + cloud hosting', ORANGE),
    ('🔔', 'Twilio SMS', 'SMS notifications for critical alerts', PURPLE),
    ('📨', 'Nodemailer', 'Email notifications and reports delivery', BLUE),
    ('⚡', 'WebSocket', 'Real-time notifications and live data updates', GREEN),
    ('📸', 'Tesseract OCR', 'Receipt and document text extraction', RED),
    ('🖥️', 'Electron', 'Native Windows/Mac desktop application', ACCENT),
    ('📱', 'PWA', 'Progressive Web App with offline support', PURPLE),
]

for i, (icon, title, desc, color) in enumerate(integrations):
    row = i % 5
    col = i // 5
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(1.8) + Inches(row * 1.0)
    add_feature_row(slide, x, y, Inches(5.8), '★', title, desc, color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 17: SECURITY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), RED)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
             '14  SECURITY', font_size=14, color=RED, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             'Security & Authentication', font_size=36, color=WHITE, bold=True)

security_features = [
    ('🔐', 'JWT Authentication', 'Access + Refresh token pair with automatic rotation', RED),
    ('🔑', 'Two-Factor Auth', 'TOTP-based 2FA (Google Authenticator / Authy)', ORANGE),
    ('🔒', 'Password Security', 'bcrypt hashing (10 salt rounds), account lockout after 5 failures', RED),
    ('🛡️', 'Encryption', 'AES-256-GCM encryption for sensitive data at rest', PURPLE),
    ('⚡', 'Rate Limiting', '100 req/15min (general), 5 req/15min (auth)', ORANGE),
    ('🪖', 'Helmet', 'Security headers (CSP, COEP, XSS protection)', BLUE),
    ('🌐', 'CORS', 'Whitelisted origin policy', ACCENT),
    ('📋', 'Audit Trail', 'Full activity logging with IP address tracking', GREEN),
    ('✅', 'Input Validation', 'express-validator on all routes', BLUE),
    ('🔌', 'Enterprise MW', 'Request IDs, API versioning, performance monitoring', ACCENT),
]

for i, (icon, title, desc, color) in enumerate(security_features):
    row = i % 5
    col = i // 5
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(1.8) + Inches(row * 1.0)
    add_feature_row(slide, x, y, Inches(5.8), '★', title, desc, color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 18: SUMMARY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
             '15  SUMMARY', font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             'Key Differentiators', font_size=36, color=WHITE, bold=True)

# Stats row
add_stat_card(slide, Inches(0.8),  Inches(1.8), Inches(1.8), Inches(1.2), '640+', 'Source Files')
add_stat_card(slide, Inches(2.9),  Inches(1.8), Inches(1.8), Inches(1.2), '150+', 'Features', GREEN)
add_stat_card(slide, Inches(5.0),  Inches(1.8), Inches(1.8), Inches(1.2), '300+', 'Endpoints', BLUE)
add_stat_card(slide, Inches(7.1),  Inches(1.8), Inches(1.8), Inches(1.2), '33K+', 'AI Code Lines', PURPLE)
add_stat_card(slide, Inches(9.2),  Inches(1.8), Inches(1.8), Inches(1.2), '38', 'AI Modules', ORANGE)
add_stat_card(slide, Inches(11.3), Inches(1.8), Inches(1.8), Inches(1.2), '0', 'External APIs', RED)

# Key differentiators
diffs = [
    ('100% Local AI', 'Every ML algorithm implemented from scratch — zero external API costs or dependencies', PURPLE),
    ('Indian Market Focus', 'INR formatting, CIBIL scoring, Indian tax rules (FY 25-26), 55+ Indian services', GREEN),
    ('Self-Learning', 'Per-user ML models that retrain automatically on data drift detection', ACCENT),
    ('3-Generation UI', 'V1 → V2 → V3 progressive enhancement — users choose their complexity level', BLUE),
    ('Enterprise Ready', 'Admin panel, role-based access, audit trails, automation engine, webhooks', ORANGE),
    ('Full Stack Platform', 'Web (React) + Desktop (Electron) + Mobile (React Native) + Cloud (Firebase)', RED),
]

for i, (title, desc, color) in enumerate(diffs):
    row = i % 3
    col = i // 3
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(3.4) + Inches(row * 1.2)
    add_feature_row(slide, x, y, Inches(5.8), '★', title, desc, color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 19: THANK YOU
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.0),
             'Thank You', font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5.5), Inches(3.2), Inches(2.3), Inches(0.04), ACCENT)

add_text_box(slide, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.6),
             'FinancialAnalyzer v2.0.0', font_size=24, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.5),
             'Circuvent Technologies', font_size=18, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.5),
             'Enterprise Financial Management  •  100% Local AI  •  Indian Market Focus',
             font_size=16, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.5),
             'Questions?',
             font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = r'c:\Users\v-hbonthada\WorkSpace-Pract\FinancialAnalyzer\docs\FinancialAnalyzer-Presentation.pptx'
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
